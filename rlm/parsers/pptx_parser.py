"""PPTX parser using python-pptx."""

from pathlib import Path

from pptx import Presentation


def parse_pptx(file_path: Path) -> str:
    """
    Extract text from a PPTX file.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text content
    """
    try:
        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())

                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)

            if len(slide_parts) > 1:  # Has content beyond slide number
                slides_text.append("\n".join(slide_parts))

        return "\n\n".join(slides_text)
    except Exception as e:
        return f"[PPTX Parse Error: {e}]"
